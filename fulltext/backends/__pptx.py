from __future__ import absolute_import

import pptx
from six import StringIO

from fulltext.util import BaseBackend
from fulltext.util import exiftool_title


class Backend(BaseBackend):

    def handle_path(self, path):
        text, p = StringIO(), pptx.Presentation(path)
        for slide in p.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text.write(run.text)
                        text.write(u'\n\n')
        return text.getvalue()

    def handle_title(self, f):
        return exiftool_title(f, self.encoding, self.encoding_errors)
